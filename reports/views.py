from django.shortcuts import render
from django.http import HttpResponse
from django.template.loader import render_to_string
from django.utils import timezone
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from weasyprint import HTML
from tracker.models import FistulaCase
from mpdsr.models import MPDSREvent
from activities.models import ActivityLog
from training.models import TrainingSession
from django.db.models import Count, Sum
from .ai_utils import generate_newsletter_narrative
from .models import MonthlyNewsletter
from .report_utils import generate_report_data
import calendar


MONTH_NAMES = {
    1:'January',2:'February',3:'March',4:'April',5:'May',6:'June',
    7:'July',8:'August',9:'September',10:'October',11:'November',12:'December'
}


def _get_pdf_context():
    now = timezone.now()
    month = now.month
    year = now.year
    data = generate_report_data(month, year)

    total_mpdsr = data['total_mpdsr']
    funded = MPDSREvent.objects.filter(action_status='FUNDED').count()

    ctx = {
        **data,
        'month_name': MONTH_NAMES[month],
        'generated_date': now.strftime('%d %b %Y'),
        'pipeline': data['pipeline'],
        'partner_data': data['partner_data'],
        'by_district': data['by_district'],
        'funded_mpdsr': funded,
        'pending_pct': round(data['pending_mpdsr'] / total_mpdsr * 100) if total_mpdsr > 0 else 0,
        'funded_pct': round(funded / total_mpdsr * 100) if total_mpdsr > 0 else 0,
        'stalled_pct': round(data['stalled_mpdsr'] / total_mpdsr * 100) if total_mpdsr > 0 else 0,
        'equity_disability': data['equity_data']['disability'],
        'equity_ethnic': data['equity_data']['ethnic_minority'],
        'equity_displaced': data['equity_data']['displaced'],
        'equity_total': sum(data['equity_data'].values()),
    }
    return ctx


def export_pdf(request):
    ctx = _get_pdf_context()
    html_string = render_to_string('one_pager.html', ctx)
    pdf = HTML(string=html_string).write_pdf()
    month_name = ctx['month_name']
    year = ctx['year']
    response = HttpResponse(pdf, content_type='application/pdf')
    response['Content-Disposition'] = f'inline; filename="CIPRB_MnE_Report_{month_name}_{year}.pdf"'
    return response


def export_ppt(request):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    UNFPA_BLUE = RGBColor(0, 158, 219)
    UNFPA_DARK = RGBColor(26, 54, 104)
    WHITE = RGBColor(255, 255, 255)
    LIGHT = RGBColor(235, 248, 255)

    now = timezone.now()
    month_name = MONTH_NAMES[now.month]
    data = generate_report_data(now.month, now.year)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = UNFPA_DARK

    tf = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    p = tf.text_frame.paragraphs[0]
    p.text = f"CIPRB M&E Monthly Report"
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE

    tf2 = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.5))
    p2 = tf2.text_frame.paragraphs[0]
    p2.text = f"{month_name} {now.year} · UNFPA Bangladesh · SRHR/RCH Programme"
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(180, 215, 235)

    # Slide 2: KPIs
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(245, 247, 250)

    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "Key Performance Indicators"
    tp.runs[0].font.size = Pt(20)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = UNFPA_DARK

    kpis = [
        ("Fistula Operated", str(data['fistula_operated']), f"of {data['fistula_goal']} target ({data['fistula_progress']}%)", UNFPA_BLUE),
        ("MPDSR Events", str(data['total_mpdsr']), f"{data['pending_mpdsr']} pending action", RGBColor(229, 62, 62)),
        ("Beneficiaries", str(data['total_beneficiaries']), "reached across activities", RGBColor(128, 90, 213)),
        ("Staff Trained", str(data['total_participants']), f"in {data['total_training_sessions']} sessions", RGBColor(56, 161, 105)),
        ("Action Rate", f"{data['action_gap_percent']}%", "MPDSR actions done", RGBColor(0, 161, 154)),
    ]

    for i, (label, value, sub, color) in enumerate(kpis):
        x = Inches(0.3 + i * 2.5)
        box = slide2.shapes.add_shape(1, x, Inches(1.0), Inches(2.3), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(226, 232, 240)

        tf = slide2.shapes.add_textbox(x + Inches(0.1), Inches(1.05), Inches(2.1), Inches(0.3))
        tf.text_frame.paragraphs[0].text = label
        tf.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        tf.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(113, 128, 150)

        tf2 = slide2.shapes.add_textbox(x + Inches(0.1), Inches(1.4), Inches(2.1), Inches(0.6))
        tf2.text_frame.paragraphs[0].text = value
        tf2.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        tf2.text_frame.paragraphs[0].runs[0].font.bold = True
        tf2.text_frame.paragraphs[0].runs[0].font.color.rgb = color

        tf3 = slide2.shapes.add_textbox(x + Inches(0.1), Inches(2.05), Inches(2.1), Inches(0.25))
        tf3.text_frame.paragraphs[0].text = sub
        tf3.text_frame.paragraphs[0].runs[0].font.size = Pt(7.5)
        tf3.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(113, 128, 150)

    # Slide 3: MPDSR + Partners
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(245, 247, 250)

    t3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    t3.text_frame.paragraphs[0].text = "MPDSR Action Status & Partner Performance"
    t3.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    t3.text_frame.paragraphs[0].runs[0].font.bold = True
    t3.text_frame.paragraphs[0].runs[0].font.color.rgb = UNFPA_DARK

    statuses = [
        ("Implemented", data['implemented_mpdsr'], RGBColor(56, 161, 105)),
        ("Pending", data['pending_mpdsr'], RGBColor(229, 62, 62)),
        ("Funded", MPDSREvent.objects.filter(action_status='FUNDED').count(), RGBColor(214, 158, 46)),
        ("Stalled", data['stalled_mpdsr'], RGBColor(113, 128, 150)),
    ]
    for i, (label, count, color) in enumerate(statuses):
        x = Inches(0.3 + i * 3.1)
        tf = slide3.shapes.add_textbox(x, Inches(1.0), Inches(2.8), Inches(1.2))
        tf.text_frame.paragraphs[0].text = str(count)
        tf.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
        tf.text_frame.paragraphs[0].runs[0].font.bold = True
        tf.text_frame.paragraphs[0].runs[0].font.color.rgb = color
        tf2 = slide3.shapes.add_textbox(x, Inches(2.0), Inches(2.8), Inches(0.4))
        tf2.text_frame.paragraphs[0].text = label
        tf2.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        tf2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(113, 128, 150)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    response = HttpResponse(output.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="CIPRB_MnE_{month_name}_{now.year}.pptx"'
    return response


def generate_newsletter(request):
    now = timezone.now()
    narrative = generate_newsletter_narrative(now.month, now.year)
    newsletter, _ = MonthlyNewsletter.objects.update_or_create(
        month=now.month,
        year=now.year,
        defaults={'content': narrative}
    )
    return render(request, 'newsletter_result.html', {'newsletter': newsletter})


def report_archive(request):
    """View all past monthly reports."""
    from .models import MonthlyReport
    reports = MonthlyReport.objects.all().order_by('-year', '-month')
    newsletters = MonthlyNewsletter.objects.all().order_by('-year', '-month')
    return render(request, 'report_archive.html', {
        'reports': reports,
        'newsletters': newsletters,
    })


# ── Design-system report views ──────────────────────────────────
from django.contrib.auth.decorators import login_required
from .live_data import get_report_json


@login_required(login_url='/accounts/login/')
def design_one_pager(request):
    return render(request, 'design/one-pager.html', {'report_json': get_report_json()})


@login_required(login_url='/accounts/login/')
def design_newsletter(request):
    return render(request, 'design/newsletter.html', {'report_json': get_report_json()})


@login_required(login_url='/accounts/login/')
def design_deck(request):
    return render(request, 'design/deck.html', {'report_json': get_report_json()})


@login_required(login_url='/accounts/login/')
def design_reports_tab(request):
    return render(request, 'design/reports-tab.html', {'report_json': get_report_json()})
